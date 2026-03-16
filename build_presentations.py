#!/usr/bin/env python3
"""
Build two research PPTX presentations:
  1. presentations/national_presentation.pptx
  2. presentations/saskatchewan_presentation.pptx
"""
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

os.makedirs("presentations", exist_ok=True)

# ── Palette ───────────────────────────────────────────────────────────────────
BG     = RGBColor(0x0F, 0x0F, 0x1A)
BG2    = RGBColor(0x1A, 0x1A, 0x2C)
CARD   = RGBColor(0x16, 0x16, 0x26)
BLUE   = RGBColor(0x4A, 0x9E, 0xFF)
RED    = RGBColor(0xE0, 0x5C, 0x5C)
ORANGE = RGBColor(0xE0, 0x88, 0x4A)
YELLOW = RGBColor(0xE0, 0xB4, 0x4A)
GREEN  = RGBColor(0x5C, 0xD6, 0x5C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FG     = RGBColor(0xE0, 0xE0, 0xE8)
FG2    = RGBColor(0x88, 0x88, 0x9A)
DARK   = RGBColor(0x07, 0x07, 0x12)
DIVIDER= RGBColor(0x20, 0x20, 0x38)

W = Inches(13.33)
H = Inches(7.5)

# ── Core helpers ──────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def box(slide, x, y, w, h, color, line_color=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=16, bold=False, color=FG,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_line(tf, text, size=14, bold=False, color=FG2,
             align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p

def bullet_box(slide, lines, x, y, w, h, size=15, header=None, header_color=BLUE):
    """Multi-line text box. lines = list of (text, color, bold, size_delta)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if header:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = header
        r.font.size = Pt(size + 1)
        r.font.bold = True
        r.font.color.rgb = header_color
    for (text, color, bold, ds) in lines:
        p = tf.paragraphs[0] if (first and not header) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size + ds)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def exposure_color(score):
    if score <= 1:   return RGBColor(0x32, 0x50, 0x32)
    if score <= 2:   return RGBColor(0x4C, 0x78, 0x28)
    if score <= 3:   return RGBColor(0x70, 0x90, 0x28)
    if score <= 4:   return RGBColor(0xA0, 0x80, 0x28)
    if score <= 5:   return RGBColor(0xC8, 0x64, 0x28)
    if score <= 6:   return RGBColor(0xE0, 0x50, 0x28)
    if score <= 7:   return RGBColor(0xE0, 0x38, 0x20)
    if score <= 8:   return RGBColor(0xE0, 0x20, 0x10)
    if score <= 9:   return RGBColor(0xEE, 0x08, 0x08)
    return                  RGBColor(0xFF, 0x00, 0x00)

# ── Slide templates ───────────────────────────────────────────────────────────
def slide_title(prs, title, subtitle="", footnote=""):
    slide = blank(prs)
    box(slide, 0, 0, W, H, BG)
    # Left accent bar
    box(slide, 0, 0, Inches(0.07), H, BLUE)
    # Top accent
    box(slide, Inches(0.07), 0, W, Inches(0.06), BLUE)
    # Bottom stripe
    box(slide, 0, H - Inches(0.06), W, Inches(0.06), RED)
    # Decorative diagonal block (top-right)
    box(slide, W - Inches(4.5), 0, Inches(4.5), Inches(2.2), DARK)

    txt(slide, title, Inches(0.7), Inches(1.6), Inches(9.5), Inches(2.8),
        size=42, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.7), Inches(4.3), Inches(10), Inches(1.0),
            size=20, color=FG2)
    if footnote:
        txt(slide, footnote, Inches(0.7), Inches(6.5), Inches(10), Inches(0.5),
            size=13, color=FG2)
    return slide

def slide_section(prs, number, title, subtitle=""):
    slide = blank(prs)
    box(slide, 0, 0, W, H, DARK)
    box(slide, 0, 0, W, Inches(0.06), BLUE)
    box(slide, 0, H - Inches(0.06), W, Inches(0.06), RED)
    # Large section number (background)
    txt(slide, f"{number:02d}", Inches(0.4), Inches(1.4), Inches(3.5), Inches(2.5),
        size=120, bold=True, color=RGBColor(0x1C, 0x1C, 0x32))
    txt(slide, f"0{number}" if number < 10 else str(number),
        Inches(0.45), Inches(1.5), Inches(3), Inches(2),
        size=90, bold=True, color=BLUE)
    txt(slide, title, Inches(0.6), Inches(3.6), Inches(11), Inches(1.6),
        size=38, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.6), Inches(5.1), Inches(11), Inches(0.8),
            size=18, color=FG2)
    return slide

def slide_content(prs, title):
    slide = blank(prs)
    box(slide, 0, 0, W, H, BG)
    box(slide, 0, 0, W, Inches(0.06), BLUE)
    txt(slide, title, Inches(0.5), Inches(0.2), Inches(12), Inches(0.72),
        size=26, bold=True, color=WHITE)
    box(slide, Inches(0.5), Inches(1.02), Inches(12.33), Inches(0.018), DIVIDER)
    return slide

def slide_two_col(prs, title, left_title="", right_title=""):
    slide = slide_content(prs, title)
    if left_title:
        txt(slide, left_title, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.5),
            size=14, bold=True, color=BLUE)
    if right_title:
        txt(slide, right_title, Inches(7.0), Inches(1.1), Inches(5.8), Inches(0.5),
            size=14, bold=True, color=BLUE)
    # Vertical divider
    box(slide, Inches(6.72), Inches(1.1), Inches(0.018), Inches(5.8), DIVIDER)
    return slide

# ── Bar chart helper ──────────────────────────────────────────────────────────
def draw_hbar(slide, rows, x, y, w, h, max_val=None, show_label=True,
              bar_h_ratio=0.55, gap_ratio=0.45):
    """rows = [(label, value, color, annotation_str)]"""
    n = len(rows)
    if max_val is None:
        max_val = max(r[1] for r in rows)
    row_h = h / n
    bar_h = row_h * bar_h_ratio
    bar_area_w = w * 0.60
    label_w = w * 0.35
    ann_w = w * 0.14

    for i, (label, val, color, ann) in enumerate(rows):
        ry = y + i * row_h + row_h * (1 - bar_h_ratio) / 2
        # Label
        txt(slide, label, x, ry + bar_h * 0.05, label_w - Inches(0.1), bar_h,
            size=11, color=FG2)
        # Track (background bar)
        bx = x + label_w
        box(slide, bx, ry, bar_area_w, bar_h, RGBColor(0x1E, 0x1E, 0x30))
        # Fill bar
        fill_w = int(bar_area_w * val / max_val)
        if fill_w > 0:
            box(slide, bx, ry, fill_w, bar_h, color)
        # Annotation
        txt(slide, ann, bx + bar_area_w + Inches(0.1), ry, ann_w, bar_h,
            size=11, bold=True, color=WHITE)

# ── Table helper ──────────────────────────────────────────────────────────────
def draw_table(slide, headers, rows, x, y, w, h, col_widths=None):
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [w / n_cols] * n_cols
    n_rows = len(rows) + 1
    row_h = h / n_rows

    # Header row
    cx = x
    for j, hdr in enumerate(headers):
        box(slide, cx, y, col_widths[j], row_h, RGBColor(0x18, 0x18, 0x2E),
            line_color=DIVIDER)
        txt(slide, hdr, cx + Inches(0.08), y + Inches(0.04),
            col_widths[j] - Inches(0.1), row_h,
            size=11, bold=True, color=BLUE)
        cx += col_widths[j]

    # Data rows
    for i, row in enumerate(rows):
        ry = y + (i + 1) * row_h
        bg_col = CARD if i % 2 == 0 else RGBColor(0x12, 0x12, 0x22)
        cx = x
        for j, cell in enumerate(row):
            box(slide, cx, ry, col_widths[j], row_h, bg_col, line_color=DIVIDER)
            cell_color = FG if j == 0 else FG2
            if isinstance(cell, tuple):
                cell_txt, cell_color = cell
            else:
                cell_txt = str(cell)
            txt(slide, cell_txt, cx + Inches(0.08), ry + Inches(0.03),
                col_widths[j] - Inches(0.1), row_h - Inches(0.04),
                size=11, color=cell_color)
            cx += col_widths[j]

# ── Callout box helper ────────────────────────────────────────────────────────
def callout(slide, lines, x, y, w, h, accent=BLUE):
    box(slide, x, y, w, h, RGBColor(0x10, 0x10, 0x25),
        line_color=RGBColor(0x28, 0x28, 0x45))
    box(slide, x, y, Inches(0.06), h, accent)
    tb = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12),
                                   w - Inches(0.28), h - Inches(0.18))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, color) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

# ═══════════════════════════════════════════════════════════════════════════════
# NATIONAL PRESENTATION
# ═══════════════════════════════════════════════════════════════════════════════
def build_national():
    prs = new_prs()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide_title(prs,
        "AI Exposure of the\nCanadian Job Market",
        "A National Assessment of 516 NOC 2021 Occupations\nUsing Open Government Data and GPT-4o Scoring",
        "516 occupations  ·  20.1M jobs  ·  GPT-4o scored  ·  March 2026")

    # ── Slide 2: Agenda ───────────────────────────────────────────────────────
    slide = slide_content(prs, "Agenda")
    sections = [
        ("01", "Background & Motivation",         "Why this study and what AI exposure means"),
        ("02", "Data Sources",                     "Five open-access Government of Canada datasets"),
        ("03", "Methodology",                      "8-step data pipeline, NOC 2021 system, GPT-4o scoring"),
        ("04", "National Findings",                "Distribution, sector breakdown, top exposures"),
        ("05", "Physical Automation",              "The other wave — robotics and dual-threat occupations"),
        ("06", "Provincial Breakdown",             "Exposure estimates across all 10 provinces"),
        ("07", "Limitations & Conclusions",        "Caveats and key takeaways"),
    ]
    for i, (num, title, sub) in enumerate(sections):
        ry = Inches(1.2) + i * Inches(0.79)
        box(slide, Inches(0.5), ry, Inches(0.52), Inches(0.52),
            BLUE if i < 3 else (RED if i < 5 else ORANGE))
        txt(slide, num, Inches(0.5), ry, Inches(0.52), Inches(0.52),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, title, Inches(1.18), ry + Inches(0.04),
            Inches(4.5), Inches(0.45), size=15, bold=True, color=FG)
        txt(slide, sub, Inches(5.9), ry + Inches(0.08),
            Inches(6.8), Inches(0.38), size=12, color=FG2)

    # ── Slide 3: Section 01 ───────────────────────────────────────────────────
    slide_section(prs, 1, "Background & Motivation",
                  "Why study AI exposure — and what does it mean?")

    # ── Slide 4: Why This Study? ──────────────────────────────────────────────
    slide = slide_content(prs, "Why This Study?")
    points = [
        "AI is reshaping knowledge work faster than any previous technology wave",
        "No existing study maps AI exposure to every NOC 2021 occupation with employment data",
        "Government of Canada open data enables a complete, quantitative, bottom-up analysis",
        "Goal: build a replicable baseline — 516 occupations, 20.1M jobs, scored 0–10",
        "Inspired by Karpathy/jobs (US BLS version) — adapted for the Canadian NOC 2021 system",
        "All data is open-access, all code is reproducible",
    ]
    for i, pt in enumerate(points):
        ry = Inches(1.25) + i * Inches(0.84)
        box(slide, Inches(0.5), ry + Inches(0.14), Inches(0.14), Inches(0.14), BLUE)
        txt(slide, pt, Inches(0.78), ry, Inches(11.5), Inches(0.75),
            size=16, color=FG)

    # ── Slide 5: What Is AI Exposure? ─────────────────────────────────────────
    slide = slide_content(prs, "What Is AI Exposure?")

    callout(slide, [
        ("AI Exposure measures COGNITIVE / DIGITAL automation — not robots or autonomous vehicles.",
         15, True, WHITE),
        ("A welder scores 2/10 even though welding robots exist.  "
         "A truck driver scores 1/10 even though AVs are coming.  "
         "Focus: does AI software change the information-processing parts of the job?",
         13, False, FG2),
    ], Inches(0.5), Inches(1.15), Inches(12.33), Inches(1.05))

    tiers = [
        ("Tier 1 · Base LLMs",
         "Writing, summarising, drafting, coding, translating, analysing\n"
         "Already changes: content creation, legal drafting, code generation",
         BLUE),
        ("Tier 2 · LLM + Tools",
         "AI with access to APIs, calendars, email, databases, CRMs, forms\n"
         "Already automates: scheduling, booking, expense processing, customer queries",
         ORANGE),
        ("Tier 3 · Multi-Agent Systems",
         "Networks of AI agents running entire workflows end-to-end\n"
         "Example: intake → contract → signing → filing with no human involvement",
         RED),
    ]
    for i, (label, body, color) in enumerate(tiers):
        cx = Inches(0.5) + i * Inches(4.16)
        box(slide, cx, Inches(2.4), Inches(3.96), Inches(4.5), CARD,
            line_color=color)
        box(slide, cx, Inches(2.4), Inches(3.96), Inches(0.07), color)
        txt(slide, label, cx + Inches(0.15), Inches(2.6), Inches(3.7),
            Inches(0.6), size=15, bold=True, color=color)
        txt(slide, body, cx + Inches(0.15), Inches(3.3), Inches(3.7),
            Inches(3.3), size=13, color=FG2)

    # ── Slide 6: Section 02 ───────────────────────────────────────────────────
    slide_section(prs, 2, "Data Sources",
                  "Five open-access Government of Canada datasets")

    # ── Slide 7: Data Sources table ───────────────────────────────────────────
    slide = slide_content(prs, "Five Open-Access Government Data Sources")
    rows = [
        ("COPS 2024–2033 Summary CSV", "ESDC",
         "516 occupations · employment 2023 · projected openings 2024–2033 · outlook"),
        ("3-Year Employment Outlooks 2025–2027", "ESDC / Job Bank",
         "~44,000 rows: regional outlook ratings + sector employment per province"),
        ("Job Bank Wage Reports (LFS 2023–2024)", "ESDC / Statistics Canada",
         "Median hourly wage by NOC unit group, 454/516 occupations, Nov 2025 update"),
        ("Table 98-10-0586-01 (2021 Census)", "Statistics Canada",
         "Median annual income by NOC unit group — fallback for 62 occupations"),
        ("NOC 2021 Structure", "Statistics Canada",
         "5-digit unit group codes, TEER levels, canonical titles, category structure"),
    ]
    cw = [Inches(3.2), Inches(2.3), Inches(6.5)]
    draw_table(slide, ["Source", "Publisher", "What we use"],
               rows, Inches(0.5), Inches(1.2), Inches(12.0), Inches(5.5), col_widths=cw)

    # ── Slide 8: NOC 2021 ─────────────────────────────────────────────────────
    slide = slide_content(prs, "The NOC 2021 Classification System")
    # Left column: explanation
    items = [
        ("National Occupational Classification 2021", True, BLUE, 15),
        ("5-digit unit groups (e.g. 21230 = Software engineers and designers)", False, FG, 14),
        ("", False, FG, 8),
        ("TEER Levels — Training, Education, Experience, Responsibilities:", True, WHITE, 14),
        ("  TEER 0 — Management", False, FG2, 13),
        ("  TEER 1 — University degree", False, FG2, 13),
        ("  TEER 2 — College / 2-yr program", False, FG2, 13),
        ("  TEER 3 — Secondary school + training", False, FG2, 13),
        ("  TEER 4 — Secondary school", False, FG2, 13),
        ("  TEER 5 — No formal education required", False, FG2, 13),
        ("", False, FG, 8),
        ("516 unit groups covered · 485 with employment data", False, FG2, 13),
        ("20.1 million jobs represented (2023 base year)", False, FG2, 13),
    ]
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, color, size) in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

    # Right: 10 categories
    cats = [
        ("Management occupations", "15", "5.5"),
        ("Business, finance & administration", "51", "5.5"),
        ("Natural & applied sciences", "66", "6.0"),
        ("Health occupations", "43", "3.5"),
        ("Education, law & social services", "52", "3.9"),
        ("Art, culture, recreation & sport", "37", "4.1"),
        ("Sales and service", "61", "2.9"),
        ("Trades, transport & equipment", "94", "2.3"),
        ("Natural resources & agriculture", "24", "2.3"),
        ("Manufacturing & utilities", "68", "2.5"),
    ]
    txt(slide, "10 Major Categories", Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.5),
        size=14, bold=True, color=BLUE)
    for i, (cat, n, avg) in enumerate(cats):
        ry = Inches(1.75) + i * Inches(0.55)
        score = float(avg)
        box(slide, Inches(7.0), ry + Inches(0.1),
            Inches(0.35), Inches(0.32), exposure_color(score))
        txt(slide, cat, Inches(7.45), ry + Inches(0.05),
            Inches(3.8), Inches(0.45), size=11, color=FG2)
        txt(slide, n, Inches(11.35), ry + Inches(0.05),
            Inches(0.5), Inches(0.45), size=11, color=FG)
        txt(slide, f"{avg}", Inches(11.9), ry + Inches(0.05),
            Inches(0.5), Inches(0.45), size=11, color=exposure_color(score))

    # ── Slide 9: Section 03 ───────────────────────────────────────────────────
    slide_section(prs, 3, "Methodology",
                  "8-step pipeline · NOC matching · GPT-4o scoring")

    # ── Slide 10: The 8-Step Pipeline ────────────────────────────────────────
    slide = slide_content(prs, "The 8-Step Data Pipeline")
    steps = [
        ("1", "build_occupations.py",       "Parse COPS CSV → 516 NOC unit groups with categories & Job Bank search URLs"),
        ("2", "build_jobbank_urls.py",       "Query Job Bank Solr API → concordance IDs for direct profile URLs (462/516 matched)"),
        ("3", "scrape_jobbank.py",           "Scrape Job Bank: canonical titles + employment requirement bullets (~3 min, resumable)"),
        ("4", "scrape_jobbank_wages.py",     "Scrape Job Bank wage reports: LFS 2023–2024 median hourly wages for 454 occupations"),
        ("5", "generate_pages.py",           "Generate one Markdown profile per occupation — structured input for the LLM scorer"),
        ("6", "make_csv_ca.py",              "Compile wages (Job Bank primary, Census fallback), employment, education, outlook → CSV"),
        ("7", "score.py",                    "Call GPT-4o API: score each occupation 0–10 + rationale. Temperature 0.2, resumable"),
        ("8", "build_site_data_ca.py",       "Merge CSV + scores + scraped Job Bank data → site/data.json (single frontend file)"),
    ]
    for i, (num, script, desc) in enumerate(steps):
        col = i % 2
        row = i // 2
        cx = Inches(0.4) + col * Inches(6.5)
        ry = Inches(1.2) + row * Inches(1.45)
        step_color = BLUE if int(num) <= 4 else RED
        box(slide, cx, ry, Inches(6.2), Inches(1.3), CARD,
            line_color=RGBColor(0x22, 0x22, 0x3A))
        box(slide, cx, ry, Inches(0.07), Inches(1.3), step_color)
        txt(slide, num, cx + Inches(0.12), ry + Inches(0.08),
            Inches(0.4), Inches(0.5), size=18, bold=True, color=step_color)
        txt(slide, script, cx + Inches(0.58), ry + Inches(0.1),
            Inches(5.4), Inches(0.45), size=13, bold=True, color=WHITE)
        txt(slide, desc, cx + Inches(0.58), ry + Inches(0.58),
            Inches(5.4), Inches(0.6), size=11, color=FG2)

    # ── Slide 11: GPT-4o Scoring ──────────────────────────────────────────────
    slide = slide_content(prs, "AI Scoring with GPT-4o")
    # Left: setup
    setup = [
        ("Model: GPT-4o (OpenAI)  ·  Temperature: 0.2", True, BLUE, 14),
        ("", False, FG, 6),
        ("Input per occupation:", True, WHITE, 14),
        ("  NOC code, TEER level, major category", False, FG2, 13),
        ("  Employment 2023 + projected openings 2024–2033", False, FG2, 13),
        ("  Labour market outlook (national modal rating)", False, FG2, 13),
        ("  Occupation title + employment requirements (scraped)", False, FG2, 13),
        ("", False, FG, 6),
        ("Output:", True, WHITE, 14),
        ('  { "exposure": N,  "rationale": "2–3 sentences" }', False, FG2, 13),
        ("", False, FG, 6),
        ("Scoring process:", True, WHITE, 14),
        ("  516 occupations scored incrementally (resumable)", False, FG2, 13),
        ("  Results cached to scores.json after each call", False, FG2, 13),
        ("  Multiple runs to clear 429 rate-limit errors", False, FG2, 13),
        ("  Final: 516/516 scored  ·  Avg (unweighted): 4.0", False, FG2, 13),
    ]
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, color, size) in setup:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

    # Right: score anchors
    anchors = [
        ("0–1  Minimal", "Roofer, landscaper, commercial diver, oil field worker", 1, GREEN),
        ("2–3  Low",     "Electrician, plumber, firefighter, dental hygienist", 2, RGBColor(0x70, 0x90, 0x28)),
        ("4–5  Moderate","Registered nurse, police officer, secondary teacher", 4, YELLOW),
        ("6–7  High",    "Accountant, journalist, HR manager, financial advisor", 6, ORANGE),
        ("8–9  Very High","Software developer, paralegal, data analyst, graphic designer", 8, RED),
        ("10  Maximum",  "Data entry clerk — routine digital, fully automatable today", 10, RGBColor(0xFF, 0x00, 0x00)),
    ]
    txt(slide, "Score Anchors", Inches(7.0), Inches(1.2), Inches(5.5), Inches(0.5),
        size=14, bold=True, color=BLUE)
    for i, (label, ex, score, color) in enumerate(anchors):
        ry = Inches(1.75) + i * Inches(0.9)
        box(slide, Inches(7.0), ry, Inches(5.7), Inches(0.82), CARD,
            line_color=RGBColor(0x22, 0x22, 0x3A))
        box(slide, Inches(7.0), ry, Inches(0.07), Inches(0.82), color)
        txt(slide, label, Inches(7.12), ry + Inches(0.05),
            Inches(5.4), Inches(0.38), size=13, bold=True, color=color)
        txt(slide, ex, Inches(7.12), ry + Inches(0.44),
            Inches(5.4), Inches(0.32), size=11, color=FG2)

    # ── Slide 12: Section 04 ─────────────────────────────────────────────────
    slide_section(prs, 4, "National Findings",
                  "Distribution, sector breakdown, and highest-exposure occupations")

    # ── Slide 13: Key Numbers ─────────────────────────────────────────────────
    slide = slide_content(prs, "Key National Statistics")
    stats = [
        ("516",    "NOC 2021 unit-group occupations\ncovered in COPS projections", BLUE),
        ("20.1M",  "Canadian jobs represented\n(2023 employment base)", BLUE),
        ("4.6/10", "Job-weighted average\nAI exposure score", YELLOW),
        ("37.6%",  "Workers in low-exposure jobs\n(score 1–3)", GREEN),
        ("38.9%",  "Workers in high-exposure jobs\n(score 6+)", RED),
        ("$590B",  "Annual payroll in high-exposure\noccupations (score 6+)", RED),
    ]
    for i, (num, desc, color) in enumerate(stats):
        col = i % 3
        row = i // 3
        cx = Inches(0.5) + col * Inches(4.22)
        ry = Inches(1.3) + row * Inches(2.7)
        box(slide, cx, ry, Inches(3.95), Inches(2.4), CARD,
            line_color=RGBColor(0x22, 0x22, 0x3A))
        box(slide, cx, ry, Inches(3.95), Inches(0.07), color)
        txt(slide, num, cx + Inches(0.2), ry + Inches(0.25),
            Inches(3.5), Inches(1.0), size=44, bold=True, color=color)
        txt(slide, desc, cx + Inches(0.2), ry + Inches(1.3),
            Inches(3.6), Inches(0.9), size=13, color=FG2)

    # ── Slide 14: Score Distribution ─────────────────────────────────────────
    slide = slide_content(prs, "Job-Weighted Score Distribution")
    dist = [
        ("Score 1", 11.3, "2.3M jobs"),
        ("Score 2", 14.4, "2.9M jobs"),
        ("Score 3", 11.9, "2.4M jobs"),
        ("Score 4", 15.3, "3.1M jobs"),
        ("Score 5",  8.2, "1.6M jobs"),
        ("Score 6", 11.5, "2.3M jobs"),
        ("Score 7", 11.2, "2.2M jobs"),
        ("Score 8", 15.1, "3.0M jobs"),
        ("Score 9",  1.0, "200K jobs"),
        ("Score 10", 0.2,  "36K jobs"),
    ]
    scores = [1, 2, 3, 4, 5, 6, 7, 8, 9, 10]
    for i, ((label, pct, ann), score) in enumerate(zip(dist, scores)):
        ry = Inches(1.22) + i * Inches(0.58)
        color = exposure_color(score)
        txt(slide, label, Inches(0.5), ry + Inches(0.04),
            Inches(1.0), Inches(0.5), size=12, color=FG2)
        # track
        box(slide, Inches(1.6), ry + Inches(0.08),
            Inches(8.0), Inches(0.38), RGBColor(0x1A, 0x1A, 0x2E))
        # fill
        fw = int(Inches(8.0) * pct / 20)  # 20% = max
        box(slide, Inches(1.6), ry + Inches(0.08), fw, Inches(0.38), color)
        txt(slide, f"{pct}%", Inches(9.7), ry + Inches(0.04),
            Inches(0.7), Inches(0.5), size=12, bold=True, color=WHITE)
        txt(slide, ann, Inches(10.5), ry + Inches(0.04),
            Inches(1.7), Inches(0.5), size=11, color=FG2)

    callout(slide, [
        ("Distribution is bimodal: high-physical-work sectors (trades, transport, NR, health) cluster at scores 1–4; "
         "knowledge-work sectors (business, finance, tech, media) cluster at 6–8. "
         "Canada's largest sector — Sales & Service — spans both.",
         12, False, FG2),
    ], Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.38), accent=YELLOW)

    # ── Slide 15: Exposure by Sector ─────────────────────────────────────────
    slide = slide_content(prs, "AI Exposure by Occupational Sector")
    sector_data = [
        ("Natural & Applied Sciences",     6.0, "66 occupations · 1.93M jobs"),
        ("Business & Finance",             5.5, "51 occupations · 3.26M jobs"),
        ("Management",                     5.5, "15 occupations · 428K jobs"),
        ("Art & Culture",                  4.1, "37 occupations · 647K jobs"),
        ("Education & Law",                3.9, "52 occupations · 2.48M jobs"),
        ("Health",                         3.5, "43 occupations · 1.62M jobs"),
        ("Sales & Service",                2.9, "61 occupations · 4.95M jobs"),
        ("Manufacturing & Utilities",      2.5, "68 occupations · 956K jobs"),
        ("Trades & Transport",             2.3, "94 occupations · 3.37M jobs"),
        ("Natural Resources & Ag",         2.3, "24 occupations · 425K jobs"),
    ]
    rows = [(lbl, val, exposure_color(val), f"{val:.1f}") for lbl, val, _ in sector_data]
    draw_hbar(slide, rows, Inches(0.5), Inches(1.15), Inches(8.5), Inches(5.9),
              max_val=7.0)
    # right side notes
    txt(slide, "Key observations", Inches(9.2), Inches(1.2), Inches(3.7), Inches(0.5),
        size=13, bold=True, color=BLUE)
    notes = [
        "Knowledge sectors score 5–6+; physical/trades sectors score 2–3",
        "Sales & Service (largest sector, 4.95M jobs) scores moderate 2.9 — but contains outliers (cashiers: 8/10)",
        "Business & Finance (3.26M jobs) averages 5.5 — highest exposure among large sectors",
        "Low AI score ≠ safe from automation — see physical robotics risk",
    ]
    tb = slide.shapes.add_textbox(Inches(9.2), Inches(1.8), Inches(3.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for note in notes:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(10)
        r = p.add_run()
        r.text = "▸  " + note
        r.font.size = Pt(12)
        r.font.color.rgb = FG2

    # ── Slide 16: Highest AI-Exposure Occupations ────────────────────────────
    slide = slide_content(prs, "Highest AI-Exposure Occupations")
    top_rows = [
        ("Data entry clerks",                          ("10/10", RED),   "36,000",  "$49K"),
        ("Receptionists",                              ("9/10", RED),  "132,000",  "$43K"),
        ("Software developers and programmers",        ("8/10", RED),  "155,700", "$100K"),
        ("Software engineers and designers",           ("8/10", RED),  "113,100", "$117K"),
        ("Administrative assistants",                  ("8/10", RED),  "233,800",  "$55K"),
        ("Administrative officers",                    ("8/10", RED),  "238,500",  "$60K"),
        ("Cashiers",                                   ("8/10", RED),  "367,300",  "$33K"),
        ("Data scientists",                            ("8/10", RED),   "36,600",  "$96K"),
        ("Financial auditors and accountants",         ("7/10", ORANGE),"247,300",  "$84K"),
        ("Information systems specialists",            ("7/10", ORANGE),"229,400",  "$96K"),
        ("Computer systems developers & programmers",  ("8/10", RED),   "45,700",  "$90K"),
    ]
    # Convert tuples properly
    formatted = []
    for occ, score_tup, jobs, pay in top_rows:
        score_str, score_color = score_tup
        formatted.append((occ, (score_str, score_color), jobs, pay))
    cw = [Inches(6.2), Inches(1.3), Inches(2.0), Inches(2.0)]
    draw_table(slide, ["Occupation", "Score", "Jobs (2023)", "Median Pay"],
               formatted, Inches(0.5), Inches(1.2), Inches(11.5), Inches(6.0), col_widths=cw)

    # ── Slide 17: Section 05 ─────────────────────────────────────────────────
    slide_section(prs, 5, "Physical Automation",
                  "The other wave — robotics, AVs, and dual-threat occupations")

    # ── Slide 18: Two Automation Waves ───────────────────────────────────────
    slide = slide_two_col(prs, "Two Distinct Automation Waves",
                          "🧠  Cognitive AI  (this study)",
                          "🦾  Physical Automation  (separate risk)")

    left_items = [
        "Language models, AI agents, multi-agent systems",
        "Attacks knowledge work, coordination, communication",
        "Score 6–10: accountants, analysts, admins, coders",
        "Mechanisms: LLMs automate drafting, analysis, scheduling",
        "Timeline: major restructuring already underway",
    ]
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(5.9), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        r = p.add_run(); r.text = "▸  " + item
        r.font.size = Pt(13); r.font.color.rgb = FG2

    right_items = [
        "Industrial robots, autonomous vehicles, precision ag equipment",
        "Attacks physical/manual work — score 1–3 occupations",
        "Trades, transport, manufacturing, agriculture, mining",
        "Brookfield (2016): 42% of Canadian jobs at high risk",
        "OECD (2023): ~27% of Canadian jobs face high risk",
    ]
    tb = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.9), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        r = p.add_run(); r.text = "▸  " + item
        r.font.size = Pt(13); r.font.color.rgb = FG2

    callout(slide, [
        ("The dual-threat framework: a complete automation picture requires BOTH dimensions. "
         "Low AI Exposure score ≠ safe from displacement — it may simply mean physical robotics "
         "is the mechanism instead.", 13, False, FG2),
    ], Inches(0.5), Inches(6.55), Inches(12.33), Inches(0.76), accent=ORANGE)

    # ── Slide 19: Dual-Threat ────────────────────────────────────────────────
    slide = slide_content(prs, "Dual-Threat Occupations")
    txt(slide, "These occupations face displacement from BOTH cognitive AI AND physical automation simultaneously:",
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5), size=13, color=FG2)

    dt_rows = [
        ("Transport truck drivers",     "1/10",  ("High", ORANGE),    "AI route optimization + AV platforms (Tesla Semi, Aurora)"),
        ("Cashiers",                    "8/10",  ("High", ORANGE),    "AI customer service + self-checkout / scan-and-go"),
        ("Agricultural workers",        "2/10",  ("High", ORANGE),    "AI crop monitoring + robotic harvesting (John Deere 8R)"),
        ("Postal workers & couriers",   "3/10",  ("High", ORANGE),    "AI route optimization + automated sorting + delivery drones"),
        ("Food processing labourers",   "2/10",  ("Very High", RED),  "Computer vision QC + fully automated processing lines"),
        ("Bank tellers",                "6/10",  ("Moderate", YELLOW),"AI banking apps + ATM / kiosk expansion"),
        ("Underground miners",          "1/10",  ("High", ORANGE),    "Autonomous haul trucks (Caterpillar, Komatsu)"),
    ]
    cw = [Inches(2.9), Inches(1.2), Inches(1.5), Inches(6.1)]
    formatted = []
    for occ, ai, (rob_str, rob_col), mech in dt_rows:
        formatted.append((occ, ai, (rob_str, rob_col), mech))
    draw_table(slide, ["Occupation", "AI Score", "Robotics", "Displacement Mechanism"],
               formatted, Inches(0.5), Inches(1.7), Inches(11.7), Inches(5.5), col_widths=cw)

    # ── Slide 20: Section 06 ─────────────────────────────────────────────────
    slide_section(prs, 6, "Provincial Breakdown",
                  "How provincial employment is estimated and what it shows")

    # ── Slide 21: Provincial Methodology ────────────────────────────────────
    slide = slide_content(prs, "How Provincial Employment Is Estimated")
    callout(slide, [
        ("COPS 2024–2033 provides national employment only. Provincial data uses the "
         "3-Year Employment Outlooks XLSX (sector-level employment per province).", 13, False, FG2),
    ], Inches(0.5), Inches(1.15), Inches(12.33), Inches(0.65), accent=BLUE)

    formula_items = [
        ("Estimation Formula:", True, WHITE, 15),
        ("", False, FG, 6),
        ("  prov_jobs = round( national_jobs × prov_cat_emp / national_cat_emp )", True, BLUE, 14),
        ("", False, FG, 6),
        ("  national_jobs     — COPS 2023 employment for this occupation (national total)", False, FG2, 13),
        ("  prov_cat_emp      — provincial employment in this occupation's major NOC category", False, FG2, 13),
        ("  national_cat_emp  — sum of all national jobs in the same major category", False, FG2, 13),
        ("", False, FG, 6),
        ("Example: SK has 18.7% of national 'Trades & Transport' workers", False, YELLOW, 13),
        ("  → Each trades occupation in SK gets 18.7% of its national count", False, FG2, 13),
        ("", False, FG, 6),
        ("Provincial avg exposure = job-weighted mean using estimated provincial counts", True, WHITE, 13),
    ]
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(7.5), Inches(4.8))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for (text, bold, color, size) in formula_items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    # Right: province table
    prov_rows = [
        ("Ontario",        "4.66", "8,227,017"),
        ("British Columbia","4.58", "2,769,384"),
        ("Quebec",         "4.58", "4,610,643"),
        ("New Brunswick",  "4.46",   "419,742"),
        ("Nova Scotia",    "4.46",   "536,697"),
        ("Alberta",        "4.45", "2,438,058"),
        ("Manitoba",       "4.39",   "700,986"),
        ("Newfoundland",   "4.38",   "254,916"),
        ("PEI",            "4.35",    "88,893"),
        ("Saskatchewan",   "4.29",   "617,463"),
    ]
    cw = [Inches(2.4), Inches(1.1), Inches(1.5)]
    draw_table(slide, ["Province", "Avg Exp.", "Est. Workers"],
               prov_rows, Inches(8.5), Inches(2.0), Inches(5.0), Inches(5.0), col_widths=cw)

    # ── Slide 22: Section 07 ─────────────────────────────────────────────────
    slide_section(prs, 7, "Limitations & Conclusions",
                  "What the data can and cannot tell us")

    # ── Slide 23: Limitations ────────────────────────────────────────────────
    slide = slide_content(prs, "Known Limitations")
    lims = [
        ("LLM opinion, not empirical measurement",
         "Scores come from GPT-4o reading text — not task-level empirical measurement. "
         "Calibration anchors reduce variance, but individual scores are tiers, not precise values.", RED),
        ("Physical automation excluded",
         "AI Exposure does not capture industrial robots, AVs, or mining systems. "
         "Low AI score ≠ safe from displacement overall.", ORANGE),
        ("Annual wages are FTE estimates",
         "Job Bank hourly medians × 2,080 hrs. Part-time workers in retail, food service, recreation "
         "typically earn less annually.", YELLOW),
        ("Employment year is 2023",
         "COPS base year. The AI landscape is moving faster than any annual survey cycle.", FG2),
        ("31 occupations lack employment data",
         "Statistics Canada suppresses figures for very small sample sizes. "
         "These occupations show at minimal size in the treemap.", FG2),
        ("National outlook is a modal proxy",
         "No national row in the 3-year outlook XLSX — we use the most common rating across "
         "13 regions, which may not reflect national conditions for regionally variable occupations.", FG2),
    ]
    for i, (title, body, color) in enumerate(lims):
        col = i % 2
        row = i // 2
        cx = Inches(0.4) + col * Inches(6.5)
        ry = Inches(1.2) + row * Inches(1.95)
        box(slide, cx, ry, Inches(6.1), Inches(1.75), CARD,
            line_color=RGBColor(0x22, 0x22, 0x3A))
        box(slide, cx, ry, Inches(0.07), Inches(1.75), color)
        txt(slide, title, cx + Inches(0.18), ry + Inches(0.1),
            Inches(5.7), Inches(0.45), size=13, bold=True, color=WHITE)
        txt(slide, body, cx + Inches(0.18), ry + Inches(0.6),
            Inches(5.7), Inches(0.98), size=11, color=FG2)

    # ── Slide 24: Conclusions ────────────────────────────────────────────────
    slide = slide_content(prs, "Conclusions")
    conclusions = [
        ("38.9% of Canadian workers (7.8M) are in high cognitive AI-exposure occupations (score 6+)",
         "$590B in annual payroll at risk of significant workflow restructuring", RED),
        ("Canada has a bimodal automation landscape",
         "Knowledge workers face cognitive AI; trades/manufacturing/ag workers face physical robotics — "
         "these two waves are distinct and require different policy responses", ORANGE),
        ("Business & Finance is the most exposed large sector",
         "3.26M jobs averaging 5.5 exposure — admin assistants, accountants, IT specialists, "
         "HR managers all face major AI-driven productivity shifts", YELLOW),
        ("The dual-threat framework matters",
         "Low AI score ≠ automation-safe. Transport, agriculture, and manufacturing workers face "
         "physical automation risk even while scoring low on cognitive AI exposure", BLUE),
        ("Open data enables reproducible analysis",
         "All five data sources are Government of Canada open access. All 8 pipeline scripts are "
         "reproducible. Full dataset and interactive tool available publicly", GREEN),
    ]
    for i, (headline, body, color) in enumerate(conclusions):
        ry = Inches(1.2) + i * Inches(1.2)
        box(slide, Inches(0.5), ry, Inches(0.45), Inches(0.45), color)
        txt(slide, str(i + 1), Inches(0.5), ry, Inches(0.45), Inches(0.45),
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, headline, Inches(1.1), ry + Inches(0.0),
            Inches(11.5), Inches(0.45), size=14, bold=True, color=WHITE)
        txt(slide, body, Inches(1.1), ry + Inches(0.48),
            Inches(11.5), Inches(0.6), size=12, color=FG2)

    prs.save("presentations/national_presentation.pptx")
    print("Saved: presentations/national_presentation.pptx")


# ═══════════════════════════════════════════════════════════════════════════════
# SASKATCHEWAN PRESENTATION
# ═══════════════════════════════════════════════════════════════════════════════
def build_saskatchewan():
    prs = new_prs()
    SK_COLOR = RGBColor(0x22, 0x55, 0x44)  # darker green for SK identity

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide_title(prs,
        "AI Exposure in\nSaskatchewan",
        "A Provincial Analysis of 617,463 Workers Across 516 Occupations",
        "Estimated from national COPS data + provincial sector employment  ·  March 2026")

    # ── Slide 2: Table of Contents ───────────────────────────────────────────
    slide = slide_content(prs, "Presentation Structure")
    sections = [
        ("01", "Background",               "National study context and what AI exposure means"),
        ("02", "The Provincial Methodology","How occupation-level SK estimates are calculated"),
        ("03", "Saskatchewan at a Glance", "Economy, employment, and sector composition"),
        ("04", "Exposure Analysis",        "Which occupations and sectors face the most AI disruption"),
        ("05", "Dual Threats",             "Occupations facing both cognitive AI and physical automation"),
        ("06", "Implications",             "What this means for SK workers, employers, and policy"),
    ]
    for i, (num, title, sub) in enumerate(sections):
        ry = Inches(1.2) + i * Inches(0.9)
        c = [BLUE, BLUE, SK_COLOR, SK_COLOR, ORANGE, RED][i]
        box(slide, Inches(0.5), ry, Inches(0.52), Inches(0.52), c)
        txt(slide, num, Inches(0.5), ry, Inches(0.52), Inches(0.52),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, title, Inches(1.18), ry + Inches(0.04),
            Inches(4.5), Inches(0.45), size=15, bold=True, color=FG)
        txt(slide, sub, Inches(5.9), ry + Inches(0.08),
            Inches(6.8), Inches(0.38), size=12, color=FG2)

    # ── Slide 3: Section 01 ───────────────────────────────────────────────────
    slide_section(prs, 1, "Background",
                  "What is AI exposure and how was this study built?")

    # ── Slide 4: National Study Context ──────────────────────────────────────
    slide = slide_content(prs, "The National Study — At a Glance")
    stats = [
        ("516",    "NOC 2021 occupations scored", BLUE),
        ("20.1M",  "Canadian jobs covered", BLUE),
        ("4.6/10", "National avg AI exposure", YELLOW),
        ("38.9%",  "Jobs in high-exposure (6+)", RED),
    ]
    for i, (num, desc, color) in enumerate(stats):
        cx = Inches(0.5) + i * Inches(3.12)
        box(slide, cx, Inches(1.2), Inches(2.85), Inches(1.8), CARD,
            line_color=RGBColor(0x22, 0x22, 0x3A))
        box(slide, cx, Inches(1.2), Inches(2.85), Inches(0.07), color)
        txt(slide, num, cx + Inches(0.15), Inches(1.35),
            Inches(2.6), Inches(0.85), size=36, bold=True, color=color)
        txt(slide, desc, cx + Inches(0.15), Inches(2.2),
            Inches(2.6), Inches(0.6), size=12, color=FG2)

    txt(slide, "Three-tier AI framework:", Inches(0.5), Inches(3.3), Inches(3.5), Inches(0.45),
        size=14, bold=True, color=WHITE)
    tiers_brief = [
        ("Tier 1 · Base LLMs", "Writing, coding, summarising, translating", BLUE),
        ("Tier 2 · LLM + Tools", "Scheduling, email, forms, CRMs, databases", ORANGE),
        ("Tier 3 · Multi-Agent", "Full workflows automated end-to-end", RED),
    ]
    for i, (t, d, c) in enumerate(tiers_brief):
        cx = Inches(0.5) + i * Inches(4.22)
        box(slide, cx, Inches(3.85), Inches(4.0), Inches(2.8), CARD,
            line_color=c)
        box(slide, cx, Inches(3.85), Inches(4.0), Inches(0.07), c)
        txt(slide, t, cx + Inches(0.15), Inches(4.05),
            Inches(3.7), Inches(0.45), size=13, bold=True, color=c)
        txt(slide, d, cx + Inches(0.15), Inches(4.6),
            Inches(3.7), Inches(1.8), size=12, color=FG2)

    callout(slide, [
        ("Score measures COGNITIVE/DIGITAL automation only — not robots, AVs, or physical automation machinery. "
         "Trades and agriculture workers face physical automation risk separately.", 12, False, FG2),
    ], Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.5), accent=YELLOW)

    # ── Slide 5: Section 02 ───────────────────────────────────────────────────
    slide_section(prs, 2, "The Provincial Methodology",
                  "How occupation-level Saskatchewan estimates are calculated")

    # ── Slide 6: Provincial Estimation Method ────────────────────────────────
    slide = slide_content(prs, "How Provincial Job Counts Are Estimated")

    callout(slide, [
        ("COPS 2024–2033 provides national employment by occupation — not by province. "
         "The 3-Year Employment Outlooks 2025–2027 XLSX provides sector-level employment per province. "
         "We combine these to estimate occupation-level provincial employment.", 13, False, FG2),
    ], Inches(0.5), Inches(1.15), Inches(12.33), Inches(0.85))

    formula = [
        ("The Formula:", True, WHITE, 16),
        ("", False, FG, 6),
        ("  prov_jobs = round( national_jobs × prov_cat_emp / national_cat_emp )", True, BLUE, 15),
        ("", False, FG, 8),
        ("Where:", True, WHITE, 14),
        ("  national_jobs    = total Canadian employment in this occupation (COPS 2023)", False, FG2, 13),
        ("  prov_cat_emp     = SK employment in this occupation's major NOC category", False, FG2, 13),
        ("  national_cat_emp = sum of all national jobs within the same NOC category", False, FG2, 13),
        ("", False, FG, 8),
        ("Worked example:", True, YELLOW, 14),
        ("  Transport truck drivers: 276,700 national jobs", False, FG2, 13),
        ("  Trades & Transport: SK has 115,592 / 3,370,000 = 3.43% of national employment", False, FG2, 13),
        ("  → SK estimated: 276,700 × 0.0343 = ~10,952 truck drivers in Saskatchewan", False, GREEN, 13),
    ]
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.15), Inches(7.8), Inches(4.8))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for (text, bold, color, size) in formula:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    # Right: validity notes
    notes_items = [
        ("Why not population share?", True, WHITE, 13),
        ("Population share ignores that SK has a disproportionately large Trades & Transport sector "
         "relative to its population — this method corrects for that.", False, FG2, 12),
        ("", False, FG, 6),
        ("Accuracy caveat:", True, YELLOW, 13),
        ("This is proportional scaling, not a survey. It estimates the occupational mix correctly "
         "at the sector level but cannot capture intra-sector differences.", False, FG2, 12),
        ("", False, FG, 6),
        ("Provincial avg exposure:", True, WHITE, 13),
        ("Computed as a job-weighted mean using the scaled SK job counts — same formula as "
         "the national average.", False, FG2, 12),
    ]
    tb = slide.shapes.add_textbox(Inches(8.5), Inches(2.15), Inches(4.4), Inches(4.8))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for (text, bold, color, size) in notes_items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

    # ── Slide 7: Section 03 ───────────────────────────────────────────────────
    slide_section(prs, 3, "Saskatchewan at a Glance",
                  "Economy, employment profile, and how SK compares nationally")

    # ── Slide 8: SK Snapshot ──────────────────────────────────────────────────
    slide = slide_content(prs, "Saskatchewan — Key Statistics")
    sk_stats = [
        ("617,463",  "Estimated total employed\nworkers in SK (2023)", BLUE),
        ("4.29/10",  "SK job-weighted\naverage AI exposure", YELLOW),
        ("4.6/10",   "National average\nAI exposure", FG2),
        ("−0.31",    "SK is below national avg\n(heavier trades/NR/ag mix)", GREEN),
    ]
    for i, (num, desc, color) in enumerate(sk_stats):
        cx = Inches(0.5) + i * Inches(3.22)
        box(slide, cx, Inches(1.2), Inches(2.95), Inches(1.85), CARD,
            line_color=RGBColor(0x22, 0x22, 0x3A))
        box(slide, cx, Inches(1.2), Inches(2.95), Inches(0.07), color)
        txt(slide, num, cx + Inches(0.15), Inches(1.35),
            Inches(2.7), Inches(0.85), size=34, bold=True, color=color)
        txt(slide, desc, cx + Inches(0.15), Inches(2.25),
            Inches(2.7), Inches(0.65), size=12, color=FG2)

    # SK sectors
    sk_sectors = [
        ("Sales & Service",           152542, 24.7, 4.26),
        ("Trades & Transport",        115592, 18.7, 2.0),
        ("Education & Law",            84973, 13.8, 4.64),
        ("Business & Finance",         82025, 13.3, 7.67),
        ("Health",                     62715, 10.2, 3.78),
        ("Natural & Applied Sci.",     41090,  6.7, 6.33),
        ("Natural Resources & Ag",     30458,  4.9, 2.32),
        ("Manufacturing & Utilities",  27208,  4.4, 2.45),
        ("Art & Culture",              15757,  2.6, 4.99),
        ("Management",                  5103,  0.8, 6.16),
    ]
    rows_table = [(s, f"{pct}%", f"{emp:,}", f"{ae}")
                  for s, emp, pct, ae in sk_sectors]
    cw = [Inches(3.2), Inches(1.0), Inches(1.6), Inches(1.2)]
    draw_table(slide, ["Sector", "Share", "Est. Jobs", "Avg Exp."],
               rows_table, Inches(0.5), Inches(3.25), Inches(7.0), Inches(4.0), col_widths=cw)

    # SK vs national comparison bars
    txt(slide, "SK vs Canada", Inches(7.8), Inches(3.25), Inches(5.0), Inches(0.5),
        size=13, bold=True, color=BLUE)
    comparisons = [
        ("Trades & Transport", 18.7, 16.7, "SK", "CA"),
        ("Natural Res & Ag",    4.9,  2.1, "SK", "CA"),
        ("Business & Finance", 13.3, 16.2, "SK", "CA"),
        ("Nat. & App. Sciences", 6.7, 9.6, "SK", "CA"),
        ("Health",             10.2, 8.1, "SK", "CA"),
    ]
    for i, (label, sk_pct, ca_pct, _, __) in enumerate(comparisons):
        ry = Inches(3.9) + i * Inches(0.65)
        txt(slide, label, Inches(7.8), ry, Inches(2.5), Inches(0.6), size=10, color=FG2)
        # SK bar
        box(slide, Inches(10.5), ry + Inches(0.06), int(Inches(2.4)*sk_pct/25), Inches(0.22),
            SK_COLOR)
        txt(slide, f"{sk_pct}% SK", Inches(10.5), ry + Inches(0.3), Inches(2.3), Inches(0.28),
            size=9, color=GREEN)
        # CA bar
        box(slide, Inches(10.5), ry + Inches(0.06), int(Inches(2.4)*ca_pct/25), Inches(0.07),
            FG2)
        txt(slide, f"{ca_pct}% CA", Inches(10.5) + int(Inches(2.4)*ca_pct/25) + Inches(0.08),
            ry, Inches(1.0), Inches(0.28), size=9, color=FG2)

    # ── Slide 9: Section 04 ───────────────────────────────────────────────────
    slide_section(prs, 4, "Exposure Analysis",
                  "Which occupations and sectors face the most AI disruption in SK?")

    # ── Slide 10: Sector Exposure ────────────────────────────────────────────
    slide = slide_content(prs, "SK AI Exposure by Sector")
    sk_sector_rows = [
        ("Business & Finance",          7.67, "82,025 SK workers · admin, finance, IT"),
        ("Natural & Applied Sciences",  6.33, "41,090 SK workers · tech, engineering"),
        ("Management",                  6.16, "5,103 SK workers"),
        ("Art & Culture",               4.99, "15,757 SK workers"),
        ("Education & Law",             4.64, "84,973 SK workers"),
        ("Sales & Service",             4.26, "152,542 SK workers — largest sector"),
        ("Health",                      3.78, "62,715 SK workers"),
        ("Manufacturing & Utilities",   2.45, "27,208 SK workers"),
        ("Natural Resources & Ag",      2.32, "30,458 SK workers"),
        ("Trades & Transport",          2.00, "115,592 SK workers — 2nd largest sector"),
    ]
    rows = [(lbl, val, exposure_color(val), f"{val:.2f}") for lbl, val, _ in sk_sector_rows]
    draw_hbar(slide, rows, Inches(0.5), Inches(1.15), Inches(9.5), Inches(5.9), max_val=9.0)

    callout(slide, [
        ("Business & Finance scores 7.67 — the highest of any sector — but employs only 13.3% of SK workers. "
         "Trades & Transport (18.7%) scores just 2.0 — pulling the provincial average down relative to ON and BC.",
         12, False, FG2),
    ], Inches(9.7), Inches(1.15), Inches(3.15), Inches(4.0), accent=YELLOW)

    # ── Slide 11: Largest SK Occupations ────────────────────────────────────
    slide = slide_content(prs, "Largest Occupations in Saskatchewan (by Est. Workers)")
    big_rows = [
        ("Retail salespersons & visual merchandisers","17,850","4/10", ("Moderate", YELLOW)),
        ("Registered nurses & psychiatric nurses",    "14,089","4/10", ("Moderate", YELLOW)),
        ("Food counter attendants & kitchen helpers", "11,795","2/10", ("Low", GREEN)),
        ("Elementary school & kindergarten teachers", "11,562","4/10", ("Moderate", YELLOW)),
        ("Nurse aides, orderlies & patient service",  "11,520","3/10", ("Low", GREEN)),
        ("Transport truck drivers",                   "10,952","1/10", ("Minimal", GREEN)),
        ("Retail & wholesale trade managers",         "10,933","6/10", ("High", ORANGE)),
        ("Cashiers",                                  "10,767","8/10", ("Very High", RED)),
        ("Managers in agriculture",                    "8,766","4/10", ("Moderate", YELLOW)),
        ("Early childhood educators & assistants",     "8,460","3/10", ("Low", GREEN)),
    ]
    cw = [Inches(4.8), Inches(1.5), Inches(1.3), Inches(1.8)]
    formatted = [(occ, jobs, score, (label, color))
                 for occ, jobs, score, (label, color) in big_rows]
    draw_table(slide, ["Occupation", "Est. SK Jobs", "AI Score", "Exposure Tier"],
               formatted, Inches(0.5), Inches(1.2), Inches(9.4), Inches(6.0), col_widths=cw)
    callout(slide, [
        ("Note: Only 3 of the 10 largest SK occupations score above 5 (cashiers 8/10, "
         "retail managers 6/10). This is why SK's avg (4.29) is below national (4.6).", 12, False, FG2),
    ], Inches(10.1), Inches(1.2), Inches(2.9), Inches(2.8), accent=BLUE)

    # ── Slide 12: High-Exposure Occupations in SK ────────────────────────────
    slide = slide_content(prs, "Highest AI-Exposure Occupations in Saskatchewan")
    txt(slide, "Occupations scoring 7+ with significant SK presence:",
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.45), size=13, color=FG2)
    high_rows = [
        ("Cashiers",                                    "10,767", ("8/10", RED),    "$33K", "Sales & Service"),
        ("Financial auditors and accountants",           "6,214", ("7/10", ORANGE), "$83K", "Business & Finance"),
        ("Administrative officers",                      "5,993", ("8/10", RED),    "$60K", "Business & Finance"),
        ("Administrative assistants",                    "5,875", ("8/10", RED),    "$54K", "Business & Finance"),
        ("Marketing & PR professionals",                 "5,870", ("8/10", RED),    "$74K", "Business & Finance"),
        ("Other customer & information services reps",   "5,596", ("8/10", RED),    "$45K", "Sales & Service"),
        ("Accounting technicians & bookkeepers",         "5,136", ("8/10", RED),    "$58K", "Business & Finance"),
        ("Information systems specialists",              "4,896", ("7/10", ORANGE), "$95K", "Natural & Applied Sci."),
        ("Accounting and related clerks",                "3,797", ("8/10", RED),    "$52K", "Business & Finance"),
        ("Receptionists",                                "3,633", ("9/10", RED),    "$43K", "Business & Finance"),
    ]
    cw = [Inches(4.0), Inches(1.4), Inches(1.2), Inches(1.2), Inches(3.5)]
    formatted = [(occ, jobs, score_tup, pay, cat)
                 for occ, jobs, score_tup, pay, cat in high_rows]
    draw_table(slide, ["Occupation", "Est. SK Jobs", "Score", "Median Pay", "Sector"],
               formatted, Inches(0.5), Inches(1.65), Inches(11.3), Inches(5.5), col_widths=cw)

    # ── Slide 13: Section 05 ─────────────────────────────────────────────────
    slide_section(prs, 5, "Dual Threats in Saskatchewan",
                  "Occupations facing both cognitive AI and physical automation")

    # ── Slide 14: Dual-Threat SK ─────────────────────────────────────────────
    slide = slide_content(prs, "Dual-Threat Occupations in Saskatchewan")
    txt(slide, "These SK occupations face displacement from BOTH cognitive AI AND physical automation:",
        Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.5), size=13, color=FG2)

    dt_sk = [
        ("Transport truck drivers",  "10,952", "1/10",  ("High", ORANGE),
         "AI route optimization + AV platforms; AV freight is a near-term SK highway reality"),
        ("Cashiers",                 "10,767", "8/10",  ("Moderate", YELLOW),
         "AI customer service + self-checkout; grocery chains already at 40%+ self-checkout"),
        ("Agricultural workers",      "~7,000", "2/10",  ("High", ORANGE),
         "AI crop monitoring + John Deere 8R autonomous tractor; SK is precision-ag ground zero"),
        ("Postal workers & couriers",  "~3,000", "3/10",  ("High", ORANGE),
         "AI route optimization + automated sorting; last-mile drone pilots in SK rural areas"),
        ("Food processing labourers",  "~2,500", "2/10",  ("Very High", RED),
         "Computer vision QC + automated processing lines; SK meat & grain processing plants"),
        ("Underground miners",         "~1,800", "1/10",  ("High", ORANGE),
         "Autonomous haul trucks (Caterpillar, Komatsu) deployed in SK potash mines"),
    ]
    cw = [Inches(2.9), Inches(1.3), Inches(1.1), Inches(1.4), Inches(5.4)]
    formatted = [(occ, jobs, ai, rob_tup, mech) for occ, jobs, ai, rob_tup, mech in dt_sk]
    draw_table(slide, ["Occupation", "SK Jobs", "AI", "Robotics", "SK-Specific Context"],
               formatted, Inches(0.5), Inches(1.65), Inches(12.1), Inches(5.5), col_widths=cw)

    # ── Slide 15: Section 06 ─────────────────────────────────────────────────
    slide_section(prs, 6, "Implications for Saskatchewan",
                  "What this means for workers, employers, and policy")

    # ── Slide 16: Implications for Workers ───────────────────────────────────
    slide = slide_content(prs, "Implications for SK Workers")
    worker_points = [
        ("~82,000 SK workers in high-exposure (6+) occupations",
         "Admin, finance, accounting, IT, and customer service workers face workflow restructuring "
         "as LLM tools and AI agents take over scheduling, drafting, and data work.",
         RED),
        ("Cashiers: the convergence point",
         "~10,767 SK cashiers score 8/10 on AI and face moderate robotics risk. "
         "This is the largest single dual-threat occupation in the province.",
         ORANGE),
        ("Trades and agriculture workers: different threat",
         "~115,592 SK trades/transport workers score just 2.0 on AI exposure — "
         "but face significant physical automation from AVs, robotic harvesters, and autonomous mining.",
         YELLOW),
        ("Healthcare workers: opportunity, not just threat",
         "62,715 SK health workers average 3.78 exposure. AI tools in healthcare extend capacity — "
         "particularly relevant given SK's ongoing nursing shortage.",
         BLUE),
        ("Reskilling gap spans education levels",
         "High-exposure occupations in SK range from cashiers (no degree) to accountants (university). "
         "A one-size-fits-all reskilling approach will not work.",
         GREEN),
    ]
    for i, (headline, body, color) in enumerate(worker_points):
        ry = Inches(1.2) + i * Inches(1.16)
        box(slide, Inches(0.5), ry, Inches(12.33), Inches(1.05), CARD,
            line_color=RGBColor(0x22, 0x22, 0x3A))
        box(slide, Inches(0.5), ry, Inches(0.07), Inches(1.05), color)
        txt(slide, headline, Inches(0.72), ry + Inches(0.06),
            Inches(11.8), Inches(0.4), size=13, bold=True, color=WHITE)
        txt(slide, body, Inches(0.72), ry + Inches(0.52),
            Inches(11.8), Inches(0.45), size=11, color=FG2)

    # ── Slide 17: Policy Implications ────────────────────────────────────────
    slide = slide_content(prs, "Implications for SK Policy & Organisations")
    policy = [
        ("Sector-specific AI upskilling — not just digital literacy",
         "Business & Finance workers (82K, avg 7.67) need AI-tool integration training. "
         "Generic 'digital skills' programs will underserve them.",
         BLUE),
        ("Agriculture & trades transition support",
         "SK's 4.9% natural resources / ag workforce is outsized vs national avg. "
         "Equipment automation transition programs should lead, not trail, deployment.",
         SK_COLOR),
        ("Healthcare: use AI to extend capacity, not reduce headcount",
         "SK faces persistent healthcare worker shortages. AI-assisted clinical admin "
         "can free up nursing and care capacity without workforce reduction.",
         GREEN),
        ("Economic diversification signal",
         "SK's lower AI exposure avg (4.29 vs 4.6) is partly structural — fewer "
         "knowledge-economy workers. Growing Business & Finance and tech sectors "
         "would raise both exposure and, typically, wage levels.",
         YELLOW),
    ]
    for i, (headline, body, color) in enumerate(policy):
        col = i % 2
        row = i // 2
        cx = Inches(0.4) + col * Inches(6.5)
        ry = Inches(1.2) + row * Inches(2.6)
        box(slide, cx, ry, Inches(6.1), Inches(2.4), CARD,
            line_color=RGBColor(0x22, 0x22, 0x3A))
        box(slide, cx, ry, Inches(0.07), Inches(2.4), color)
        txt(slide, headline, cx + Inches(0.18), ry + Inches(0.1),
            Inches(5.7), Inches(0.5), size=13, bold=True, color=WHITE)
        txt(slide, body, cx + Inches(0.18), ry + Inches(0.65),
            Inches(5.7), Inches(1.6), size=12, color=FG2)

    # ── Slide 18: Conclusions ────────────────────────────────────────────────
    slide = slide_content(prs, "Conclusions — Saskatchewan")
    conclusions_sk = [
        ("SK avg exposure 4.29/10 — below national 4.6",
         "Driven by a heavier trades, transport, and natural resources mix vs ON/BC. "
         "SK has proportionally more workers in physically-demanding, low-AI-exposure sectors.", FG2),
        ("Business & Finance: SK's most exposed large sector",
         "82,025 workers averaging 7.67/10. Admin assistants, accountants, receptionists, "
         "and IT specialists face the most direct AI workflow disruption.", RED),
        ("~82,000 SK workers in high-exposure (6+) occupations",
         "Concentrated in admin, finance, sales coordination, and IT. Wage range is wide: "
         "$33K (cashiers) to $95K (information systems specialists).", ORANGE),
        ("Dual threats are a SK-specific reality",
         "Agriculture (autonomous harvesters), transport (AV trucks), and potash mining "
         "(autonomous haul trucks) mean SK faces physical automation alongside cognitive AI.", YELLOW),
        ("Full provincial dataset available",
         "Interactive treemap and occupation-level data at: /province/SK.html", BLUE),
    ]
    for i, (headline, body, color) in enumerate(conclusions_sk):
        ry = Inches(1.2) + i * Inches(1.2)
        box(slide, Inches(0.5), ry, Inches(0.45), Inches(0.45), color)
        txt(slide, str(i + 1), Inches(0.5), ry, Inches(0.45), Inches(0.45),
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, headline, Inches(1.1), ry + Inches(0.0),
            Inches(11.5), Inches(0.45), size=14, bold=True, color=WHITE)
        txt(slide, body, Inches(1.1), ry + Inches(0.48),
            Inches(11.5), Inches(0.6), size=12, color=FG2)

    prs.save("presentations/saskatchewan_presentation.pptx")
    print("Saved: presentations/saskatchewan_presentation.pptx")


if __name__ == "__main__":
    build_national()
    build_saskatchewan()
    print("\nDone! Both presentations saved to presentations/")
